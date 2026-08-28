# -*- coding: utf-8 -*-
"""
build_pptx.py —— 把结构化的开题内容生成为开题报告 PPT。

用法:
    python build_pptx.py content.json [输出.pptx]

- 若已安装 python-pptx，生成真正的 .pptx（16:9，含封面/目录/正文/致谢，带演讲者备注）。
- 若未安装 python-pptx，自动降级：生成同名 .md（Marp 幻灯格式），并提示安装方法。

content.json 结构见 references/ppt_structure.md，示例见 scripts/content.example.json。
"""

import sys
import os
import json

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# ---------- 主题配色 ----------
PRIMARY = (0x1F, 0x38, 0x64)   # 深蓝：标题
ACCENT = (0x2E, 0x75, 0xB6)    # 亮蓝：装饰条
TEXT = (0x33, 0x33, 0x33)      # 正文
LIGHT = (0xF2, 0xF5, 0xFA)     # 浅底


def load_content(path):
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    data.setdefault("meta", {})
    data.setdefault("slides", [])
    return data


# =========================================================================
# 路径 A：python-pptx 生成真正的 .pptx
# =========================================================================

def build_with_pptx(data, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    def rgb(t):
        return RGBColor(*t)

    def add_slide():
        return prs.slides.add_slide(BLANK)

    def add_rect(slide, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def add_text(slide, x, y, w, h, text, size, color, bold=False,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = rgb(color)
        r.font.name = "微软雅黑"
        return tb

    def set_notes(slide, notes):
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # ---- 封面 ----
    meta = data["meta"]
    s = add_slide()
    add_rect(s, 0, 0, SW, SH, LIGHT)
    add_rect(s, 0, Inches(2.7), SW, Inches(0.08), ACCENT)
    add_text(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.4),
             meta.get("title", "开题报告"), 40, PRIMARY, bold=True,
             align=PP_ALIGN.CENTER)
    if meta.get("subtitle"):
        add_text(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.8),
                 meta["subtitle"], 22, ACCENT, align=PP_ALIGN.CENTER)
    info = []
    for label, key in [("汇报人", "author"), ("指导教师", "advisor"),
                       ("单位", "institution"), ("日期", "date")]:
        if meta.get(key):
            info.append(f"{label}：{meta[key]}")
    if info:
        add_text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.0),
                 "\n".join(info), 20, TEXT, align=PP_ALIGN.CENTER)
    set_notes(s, meta.get("notes", ""))

    # ---- 目录 ----
    titles = [sl.get("title", "") for sl in data["slides"]]
    if titles:
        s = add_slide()
        add_rect(s, 0, 0, SW, Inches(1.2), PRIMARY)
        add_text(s, Inches(0.6), Inches(0.28), Inches(12), Inches(0.8),
                 "汇报提纲", 30, (255, 255, 255), bold=True)
        lines = "\n".join(f"{i+1}.  {t}" for i, t in enumerate(titles))
        add_text(s, Inches(1.2), Inches(1.6), Inches(11), Inches(5.4),
                 lines, 22, TEXT, anchor=MSO_ANCHOR.TOP)

    # ---- 正文 ----
    total = len(data["slides"])
    for idx, sl in enumerate(data["slides"], 1):
        s = add_slide()
        # 顶部标题条
        add_rect(s, 0, 0, SW, Inches(1.2), PRIMARY)
        add_rect(s, 0, Inches(1.2), SW, Inches(0.06), ACCENT)
        add_text(s, Inches(0.6), Inches(0.28), Inches(11.5), Inches(0.8),
                 f"{idx}. {sl.get('title','')}", 28, (255, 255, 255), bold=True)
        # 正文要点
        bullets = sl.get("bullets", [])
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6),
                                  Inches(11.7), Inches(5.2))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for b in bullets:
            if isinstance(b, (list, tuple)):
                txt, lvl = b[0], (b[1] if len(b) > 1 else 0)
            elif isinstance(b, dict):
                txt, lvl = b.get("text", ""), b.get("level", 0)
            else:
                txt, lvl = str(b), 0
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = min(int(lvl), 4)
            mark = "• " if lvl == 0 else "– "
            run = p.add_run(); run.text = mark + txt
            run.font.size = Pt(22 - min(int(lvl), 2) * 3)
            run.font.color.rgb = rgb(PRIMARY if lvl == 0 else TEXT)
            run.font.bold = (lvl == 0)
            run.font.name = "微软雅黑"
            p.space_after = Pt(8)
        # 页脚页码
        add_text(s, Inches(11.8), Inches(7.0), Inches(1.4), Inches(0.4),
                 f"{idx} / {total}", 12, TEXT, align=PP_ALIGN.RIGHT)
        set_notes(s, sl.get("notes", ""))

    # ---- 致谢 ----
    s = add_slide()
    add_rect(s, 0, 0, SW, SH, PRIMARY)
    add_text(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2),
             data.get("thanks", "恳请各位专家批评指正！"), 36,
             (255, 255, 255), bold=True, align=PP_ALIGN.CENTER)

    prs.save(out_path)
    print(f"[OK] 已生成 PPT：{out_path}（共 {len(prs.slides._sldIdLst)} 页）")


# =========================================================================
# 路径 B：降级为 Marp 幻灯 Markdown
# =========================================================================

def build_marp(data, out_md):
    meta = data["meta"]
    L = ["---", "marp: true", "theme: default", "paginate: true", "---", ""]
    # 封面
    L.append(f"# {meta.get('title','开题报告')}")
    if meta.get("subtitle"):
        L.append(f"### {meta['subtitle']}")
    L.append("")
    for label, key in [("汇报人", "author"), ("指导教师", "advisor"),
                       ("单位", "institution"), ("日期", "date")]:
        if meta.get(key):
            L.append(f"**{label}**：{meta[key]}")
    L.append("")
    # 目录
    titles = [sl.get("title", "") for sl in data["slides"]]
    if titles:
        L += ["---", "", "## 汇报提纲", ""]
        L += [f"{i+1}. {t}" for i, t in enumerate(titles)]
        L.append("")
    # 正文
    for idx, sl in enumerate(data["slides"], 1):
        L += ["---", "", f"## {idx}. {sl.get('title','')}", ""]
        for b in sl.get("bullets", []):
            if isinstance(b, (list, tuple)):
                txt, lvl = b[0], (b[1] if len(b) > 1 else 0)
            elif isinstance(b, dict):
                txt, lvl = b.get("text", ""), b.get("level", 0)
            else:
                txt, lvl = str(b), 0
            L.append("  " * int(lvl) + f"- {txt}")
        L.append("")
        if sl.get("notes"):
            L.append(f"<!-- 备注: {sl['notes']} -->")
            L.append("")
    # 致谢
    L += ["---", "", f"# {data.get('thanks', '恳请各位专家批评指正！')}", ""]
    with open(out_md, "w", encoding="utf-8") as f:
        f.write("\n".join(L))
    print(f"[OK] 已生成 Marp 幻灯：{out_md}")
    print("     渲染为 PPTX/PDF： npx @marp-team/marp-cli " +
          os.path.basename(out_md) + " --pptx")


# =========================================================================

def main():
    if len(sys.argv) < 2:
        print("用法: python build_pptx.py content.json [输出.pptx]")
        sys.exit(1)
    src = sys.argv[1]
    data = load_content(src)
    out = sys.argv[2] if len(sys.argv) > 2 else "开题报告PPT.pptx"

    try:
        import pptx  # noqa: F401
    except ImportError:
        print("[提示] 未检测到 python-pptx，无法直接生成 .pptx。")
        print("       安装后可生成：pip install python-pptx")
        print("       现降级生成 Marp 幻灯 Markdown 作为替代。")
        md = os.path.splitext(out)[0] + ".md"
        build_marp(data, md)
        return

    build_with_pptx(data, out)


if __name__ == "__main__":
    main()
